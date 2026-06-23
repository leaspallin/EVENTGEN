"""Moteur de génération : remplit le VRAI template ALL IN puis exporte PPTX + PDF."""
import copy, io, os, shutil, subprocess, tempfile
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from pptx.text.text import _Paragraph
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = os.path.join(HERE, "template", "template.pptx")

# ---- mapping slide(0-based) -> shape_id ----
S = {
    "cover_logo_client": (0, 1026),
    "date":            (1, 18),
    "programme":       (2, 6),
    "espace_titre":    (3, 4),
    "espace_specs":    (3, 17),
    "espace_photo":    (3, 3),
    "devis_tarif":     (11, 26),
    "devis_incluant":  (11, 5),
    "devis_table":     (11, 3),
}

def _shape(slide, sid):
    for sh in slide.shapes:
        if sh.shape_id == sid:
            return sh
    raise KeyError(f"shape {sid} introuvable")

def _set_run_line(shape, text):
    p = shape.text_frame.paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        for r in p.runs[1:]:
            r.text = ""
    else:
        p.text = text

def _set_two_lines(shape, l1, l2):
    paras = shape.text_frame.paragraphs
    for p, val in zip(paras, [l1, l2]):
        if p.runs:
            p.runs[0].text = val
            for r in p.runs[1:]:
                r.text = ""
        else:
            p.text = val


def _set_room_name(shape, room):
    """Le titre d'espace est un paragraphe a 2 runs (ESPACE / BERCY).
       On remplace uniquement le run du nom de salle, on garde 'ESPACE'."""
    runs=[r for p in shape.text_frame.paragraphs for r in p.runs]
    target=None
    for r in runs:
        u=r.text.upper()
        if 'BERCY' in u or 'TORONTO' in u:
            target=r; break
    if target is None and len(runs)>1:
        target=runs[-1]
    if target is not None:
        target.text=room

def _set_header_block(shape, header, items, bullet_index=1):
    """1er paragraphe = header ; puis 1 paragraphe à puce cloné par item.
       Garantit des puces homogènes (pas de double tiret)."""
    tf = shape.text_frame
    txBody = tf._txBody
    paras = tf.paragraphs
    # header sur le paragraphe 0
    p0 = paras[0]
    if p0.runs:
        p0.runs[0].text = header
        for r in p0.runs[1:]:
            r.text = ""
    else:
        p0.text = header
    # paragraphe-modèle à puce
    idx = bullet_index if len(paras) > bullet_index else (1 if len(paras) > 1 else 0)
    bullet_proto = copy.deepcopy(paras[idx]._p)
    # supprimer tous les paragraphes après le header
    for p in list(paras[1:]):
        txBody.remove(p._p)
    # ajouter un paragraphe par item
    for it in items:
        new = copy.deepcopy(bullet_proto)
        txBody.append(new)
        np = _Paragraph(new, tf)
        if np.runs:
            np.runs[0].text = it
            for r in np.runs[1:]:
                r.text = ""
        else:
            np.text = it

def _set_cell(table, r, c, text):
    p = table.cell(r, c).text_frame.paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        for rr in p.runs[1:]:
            rr.text = ""
    else:
        p.text = text

def _replace_picture(slide, shape, image_bytes):
    """Remplace le blob de l'image en conservant cadre/position ; ajuste le ratio."""
    bio = io.BytesIO(image_bytes)
    image_part, rId = slide.part.get_or_add_image_part(bio)
    pic = shape._element
    for blip in pic.iter(qn('a:blip')):
        blip.set(qn('r:embed'), rId)
        break
    # ajuster pour préserver le ratio dans le cadre d'origine (centré)
    try:
        im = Image.open(io.BytesIO(image_bytes))
        iw, ih = im.size
        bw, bh = shape.width, shape.height
        bx, by = shape.left, shape.top
        cx, cy = bx + bw / 2, by + bh / 2
        scale = min(bw / iw, bh / ih)
        nw, nh = int(iw * scale), int(ih * scale)
        shape.width, shape.height = nw, nh
        shape.left, shape.top = int(cx - nw / 2), int(cy - nh / 2)
    except Exception:
        pass

ESPACE_PRESETS = {
    "bercy":   ("ESPACE", "BERCY", [
        "1 vidéoprojecteur laser Phosphor WUXGA",
        "2 micros HF main",
        "1 écran 32 pouces à l'entrée de l'espace pour le logo entreprise",
        "Équipements clickshare CX-30 (conférence et collaboration sans fil)",
        "Matrice audio numérique",
    ]),
    "toronto": ("ESPACE", "TORONTO", [
        "1 écran LED 5m x 2,5m",
        "2 micros HF main",
        "1 écran 32 pouces à l'entrée de l'espace pour le logo entreprise",
        "Équipements clickshare CX-30 (conférence et collaboration sans fil)",
        "Matrice audio numérique",
    ]),
    "combo":   ("ESPACE", "BERCY + TORONTO", [
        "1 écran LED 5m x 2,5m + vidéoprojecteur laser WUXGA",
        "2 micros HF main",
        "1 écran 32 pouces à l'entrée de l'espace pour le logo entreprise",
        "Équipements clickshare CX-30 (conférence et collaboration sans fil)",
        "Matrice audio numérique",
    ]),
}

def compute_devis(lines, pax, remise_pct, tva_pct):
    net = sum((l["montant"] * pax if l["type"] == "pax" else l["montant"]) for l in lines)
    rem = remise_pct / 100.0
    brut = net / (1 - rem) if rem < 1 else net
    tva = net * tva_pct / 100.0
    return {
        "net": net, "brut": brut, "tva": tva, "ttc": net + tva,
        "pp_net": net / pax if pax else 0,
        "pp_brut": brut / pax if pax else 0,
    }

def _eur(n):
    return f"{n:,.2f} €".replace(",", " ").replace(".", ",").replace(" ", "\u00a0", 1) \
           .replace(",", "§").replace(".", ",").replace("§", ",") if False else \
           f"{n:,.2f}\u00a0€".replace(",", " ").replace(".", ",")

def build(data):
    prs = Presentation(TEMPLATE)
    sl = list(prs.slides)

    # logo client (couverture)
    if data.get("logo_bytes"):
        _replace_picture(sl[0], _shape(sl[0], S["cover_logo_client"][1]), data["logo_bytes"])

    # date
    _set_run_line(_shape(sl[1], S["date"][1]), data["dates"])

    # programme
    prog_items = [l for l in data["programme"] if l.strip()]
    _set_header_block(_shape(sl[2], S["programme"][1]),
                      f"Pour {data['pax']} personnes :", prog_items)

    # espace
    t1, t2, specs = ESPACE_PRESETS[data["espace"]]
    _set_room_name(_shape(sl[3], S["espace_titre"][1]), t2)
    _set_header_block(_shape(sl[3], S["espace_specs"][1]),
                      "Matériel technique mis à disposition :", specs)
    if data.get("espace_photo_bytes"):
        _replace_picture(sl[3], _shape(sl[3], S["espace_photo"][1]), data["espace_photo_bytes"])

    # devis
    d = compute_devis(data["devis_lines"], data["pax"], data["remise"], data["tva"])
    _set_run_line(_shape(sl[11], S["devis_tarif"][1]),
                  f"TARIF PAR PERSONNE : {_eur(d['pp_net'])} HT")
    incl = [l["label"] for l in data["devis_lines"] if l["label"].strip()]
    _set_header_block(_shape(sl[11], S["devis_incluant"][1]), "Incluant :", incl)
    tbl = _shape(sl[11], S["devis_table"][1]).table
    _set_cell(tbl, 0, 1, _eur(d["net"]))
    _set_cell(tbl, 1, 1, _eur(d["tva"]))
    _set_cell(tbl, 2, 1, _eur(d["ttc"]))

    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io.getvalue(), d

def _soffice_bin():
    for b in ("soffice", "libreoffice"):
        if shutil.which(b):
            return b
    return None

def pdf_available():
    return _soffice_bin() is not None

def to_pdf(pptx_bytes):
    binr = _soffice_bin()
    if not binr:
        raise FileNotFoundError("Convertisseur PDF (LibreOffice) indisponible")
    with tempfile.TemporaryDirectory() as tmp:
        src = os.path.join(tmp, "deck.pptx")
        with open(src, "wb") as f:
            f.write(pptx_bytes)
        subprocess.run([binr, "--headless", "--convert-to", "pdf",
                        "--outdir", tmp, src],
                       check=True, timeout=180,
                       stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        pdf = os.path.join(tmp, "deck.pdf")
        with open(pdf, "rb") as f:
            return f.read()
